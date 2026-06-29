from __future__ import annotations

import csv
import io
import json
import re
from dataclasses import dataclass
from html import unescape
from typing import Any

from bs4 import BeautifulSoup


@dataclass
class ExtractedDocument:
    text: str
    provenance: list[dict[str, Any]]
    warnings: list[dict[str, Any]]


def _unit(text: str, **metadata: Any) -> dict[str, Any]:
    return {"text": text[:240], **metadata}


def extract_document(filename: str, data: bytes, file_type: str) -> ExtractedDocument:
    ext = file_type.lower().lstrip(".")
    if not data:
        return ExtractedDocument("", [], [{"code": "empty_file", "message": "File is empty."}])

    try:
        if ext == "pdf":
            return _extract_pdf(data)
        if ext == "docx":
            return _extract_docx(data)
        if ext == "pptx":
            return _extract_pptx(data)
        if ext == "xlsx":
            return _extract_xlsx(data)
        if ext == "csv":
            return _extract_csv(data)
        if ext in {"html", "htm"}:
            return _extract_html(data)
        if ext == "json":
            parsed = json.loads(_decode_text(data))
            text = json.dumps(parsed, indent=2, ensure_ascii=False)
            return ExtractedDocument(text, [_unit(text, section="json")], [])
        if ext == "xml":
            text = _strip_xml_noise(_decode_text(data))
            return ExtractedDocument(text, [_unit(text, section="xml")], [])
        return _extract_plain_text(data)
    except Exception as exc:
        return ExtractedDocument(
            "",
            [],
            [{"code": "unreadable_file", "message": f"Could not extract text from {filename}: {exc}"}],
        )


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_plain_text(data: bytes) -> ExtractedDocument:
    text = _decode_text(data)
    return ExtractedDocument(text, [_unit(text, section="text", start_offset=0, end_offset=len(text))], [])


def _extract_pdf(data: bytes) -> ExtractedDocument:
    import fitz

    doc = fitz.open(stream=data, filetype="pdf")
    parts: list[str] = []
    provenance: list[dict[str, Any]] = []
    for index, page in enumerate(doc, start=1):
        page_text = page.get_text("text")
        if page_text.strip():
            parts.append(f"\n\n[Page {index}]\n{page_text}")
            provenance.append(_unit(page_text, page_number=index))
    warnings = []
    if not "".join(parts).strip():
        warnings.append(
            {
                "code": "scanned_pdf_no_text",
                "message": "PDF has no extractable text. OCR can be added to this pipeline later.",
            }
        )
    return ExtractedDocument("\n".join(parts).strip(), provenance, warnings)


def _extract_docx(data: bytes) -> ExtractedDocument:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(data))
    parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    text = "\n".join(parts)
    return ExtractedDocument(text, [_unit(value, paragraph=i + 1) for i, value in enumerate(parts)], [])


def _extract_pptx(data: bytes) -> ExtractedDocument:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    provenance: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            text = "\n".join(slide_parts)
            parts.append(f"[Slide {slide_number}]\n{text}")
            provenance.append(_unit(text, slide_number=slide_number))
    return ExtractedDocument("\n\n".join(parts), provenance, [])


def _extract_xlsx(data: bytes) -> ExtractedDocument:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts: list[str] = []
    provenance: list[dict[str, Any]] = []
    for sheet in workbook.worksheets:
        rows = []
        for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = [str(value) if value is not None else "" for value in row]
            if any(value.strip() for value in values):
                rows.append(" | ".join(values))
                provenance.append(_unit(" | ".join(values), sheet_name=sheet.title, row=row_number))
        if rows:
            parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return ExtractedDocument("\n\n".join(parts), provenance, [])


def _extract_csv(data: bytes) -> ExtractedDocument:
    text = _decode_text(data)
    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(row) for row in reader]
    return ExtractedDocument("\n".join(rows), [_unit(row, row=i + 1) for i, row in enumerate(rows)], [])


def _extract_html(data: bytes) -> ExtractedDocument:
    soup = BeautifulSoup(_decode_text(data), "html.parser")
    for element in soup(["script", "style", "nav", "footer", "aside"]):
        element.decompose()
    title = soup.title.string.strip() if soup.title and soup.title.string else ""
    text = soup.get_text("\n")
    lines = [unescape(line.strip()) for line in text.splitlines() if line.strip()]
    body = "\n".join(lines)
    if title and title not in body:
        body = f"{title}\n{body}"
    return ExtractedDocument(body, [_unit(body, section="html")], [])


def _strip_xml_noise(text: str) -> str:
    text = re.sub(r"<\?xml[^>]*>", "", text)
    return re.sub(r">\s+<", ">\n<", text).strip()
